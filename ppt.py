import sys
try:
    from pptx import Presentation
except ImportError:
    print("请先在终端安装 python-pptx 库：")
    print("pip install python-pptx")
    sys.exit()

def create_presentation():
    # 创建一个新的PPT对象
    prs = Presentation()

    # 幻灯片 1：封面页 (布局索引 0: 标题幻灯片)
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "人工智能 (AI)：探索定义、应用与未来趋势"
    subtitle.text = "重塑世界的新引擎\n汇报人：[您的姓名/团队名称]\n日期：2026年[月]日"

    # 幻灯片 2：目录页 (布局索引 1: 标题和内容)
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "目录 CONTENTS"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "什么是人工智能？ (AI的定义与分类)"
    tf.add_paragraph().text = "AI 的核心技术底座"
    tf.add_paragraph().text = "AI 的主要应用领域 (医疗、交通、金融、制造)"
    tf.add_paragraph().text = "AI 的未来趋势与挑战"
    tf.add_paragraph().text = "总结与展望"

    # 幻灯片 3：什么是人工智能？
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "探寻本源：什么是人工智能？"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "核心定义："
    p = tf.add_paragraph()
    p.text = "人工智能（Artificial Intelligence）是计算机科学的一个分支。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "它企图了解智能的实质，并生产出一种新的能以人类智能相似方式做出反应的智能机器。"
    p.level = 1
    tf.add_paragraph().text = "AI 的两个阶段/分类："
    p = tf.add_paragraph()
    p.text = "弱人工智能 (ANI)：擅长单个特定任务（如语音助手、下棋软件）。我们目前正处于这一阶段。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "强人工智能 (AGI)：具备与人类同等甚至超越人类的全面认知能力（未来的发展目标）。"
    p.level = 1

    # 幻灯片 4：核心技术底座
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "支撑AI运作的关键技术"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "机器学习 (Machine Learning)：让计算机在没有明确编程的情况下，从数据中学习规律。"
    tf.add_paragraph().text = "深度学习 (Deep Learning)：基于人工神经网络，能够处理海量且复杂的非结构化数据（如图像、音频）。"
    tf.add_paragraph().text = "自然语言处理 (NLP)：让机器理解、解释和生成人类语言（例如目前的ChatGPT、文心一言等大语言模型）。"
    tf.add_paragraph().text = "计算机视觉 (CV)：使计算机能够“看”并理解数字图像或视频（如人脸识别、医学影像分析）。"

    # 幻灯片 5：应用领域一 —— 智慧医疗
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "AI 赋能医疗：让生命更有保障"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "辅助诊疗：AI 识图技术快速分析X光、CT影像，准确率比肩资深医师，有效降低误诊率。"
    tf.add_paragraph().text = "新药研发：利用大模型预测蛋白质结构，将靶点发现和药物合成的周期从数年缩短至数月。"
    tf.add_paragraph().text = "个性化健康管理：可穿戴设备结合AI算法，实时监测心率、血糖，提供疾病预警。"

    # 幻灯片 6：应用领域二 —— 智能交通
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "AI 改变出行：更安全、更高效"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "自动驾驶汽车：通过激光雷达与机器视觉，实现环境感知、路径规划与自动避障。"
    tf.add_paragraph().text = "智慧交通调度：城市交通大脑实时分析路况，动态调整红绿灯时长，缓解城市拥堵。"
    tf.add_paragraph().text = "智能物流：无人驾驶卡车、无人机配送与仓储机器人协同，大幅提升物流流转效率。"

    # 幻灯片 7：应用领域三 —— 金融科技
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "AI 驱动金融：精准、安全与智能化"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "智能风控与反欺诈：毫秒级分析交易行为，精准识别盗刷、洗钱等异常活动。"
    tf.add_paragraph().text = "量化交易与投资：AI 算法分析全球市场数据与新闻情绪，捕捉稍纵即逝的投资机会。"
    tf.add_paragraph().text = "智能客服系统：7x24小时全天候响应，通过语音和自然语言处理解决大部分基础客户查询。"

    # 幻灯片 8：应用领域四 —— 智能制造
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "AI 无处不在：从工厂到家庭"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "工业 4.0 (智能制造)："
    p = tf.add_paragraph()
    p.text = "工业机器人实现高精度组装。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "设备的预测性维护，在故障发生前进行修复，减少停机损失。"
    p.level = 1
    tf.add_paragraph().text = "智慧生活："
    p = tf.add_paragraph()
    p.text = "全屋智能家居（智能音箱控制家电、AI灯光调节）。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "个性化推荐系统（懂你喜好的短视频、音乐和电商购物推送）。"
    p.level = 1

    # 幻灯片 9：未来趋势与面临的挑战
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "AI 的明天：机遇与挑战并存"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "未来趋势："
    p = tf.add_paragraph()
    p.text = "端侧大模型：AI 将直接部署在手机、PC端，离线可用且更保护隐私（Edge AI）。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "多模态融合：未来的AI能同时处理文本、图像、视频和音频，更像真实人类。"
    p.level = 1
    tf.add_paragraph().text = "关键挑战："
    p = tf.add_paragraph()
    p.text = "数据隐私与安全：如何防止数据泄露与滥用？"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "伦理与就业：AI产生的版权归属？如何解决AI替代部分基础岗位的问题？"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "算法黑盒：如何让深度学习的决策过程更具“可解释性”？"
    p.level = 1

    # 幻灯片 10：结语与致谢
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "结语与致谢 (Q&A)"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "核心观点："
    p = tf.add_paragraph()
    p.text = "人工智能不是为了取代人类，而是作为强大的工具赋能人类。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "拥抱变化，终身学习，与AI共同进化是我们面对未来的最佳态度。"
    p.level = 1
    tf.add_paragraph().text = "感谢聆听！"
    tf.add_paragraph().text = "Q & A / 提问环节"

    # 保存文件
    prs.save('AI_Presentation.pptx')
    print("成功生成 PPT 文件：AI_Presentation.pptx")
    print("您现在可以直接用 PowerPoint 打开它并应用您喜欢的设计主题了！")

if __name__ == '__main__':
    create_presentation()